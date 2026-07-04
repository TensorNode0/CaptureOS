"""Document exports: markdown -> .docx, cost volume -> .xlsx, deck -> .pptx,
capability one-pager -> .docx, and full-package .zip — all in memory."""
import io
import re
import zipfile
from datetime import datetime, timezone

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt


def safe_filename(name, ext):
    slug = re.sub(r"[^A-Za-z0-9._ -]+", "", name or "document").strip()
    slug = re.sub(r"\s+", "_", slug)[:80] or "document"
    return f"{slug}.{ext}"


# ─────────────────────────── markdown -> docx ────────────────────────────────

_BOLD_RE = re.compile(r"(\*\*.+?\*\*|\*[^*]+?\*)")


def _add_runs(paragraph, text):
    """Render **bold** / *italic* inline markdown into docx runs."""
    pos = 0
    for m in _BOLD_RE.finditer(text):
        if m.start() > pos:
            paragraph.add_run(text[pos:m.start()])
        token = m.group(0)
        if token.startswith("**"):
            run = paragraph.add_run(token[2:-2])
            run.bold = True
        else:
            run = paragraph.add_run(token[1:-1])
            run.italic = True
        pos = m.end()
    if pos < len(text):
        paragraph.add_run(text[pos:])


def _is_table_sep(line):
    body = line.strip().strip("|").strip()
    return bool(body) and all(c in "-: " for c in body)


def _split_row(line):
    return [c.strip() for c in line.strip().strip("|").split("|")]


def _add_md_table(doc, header, rows):
    table = doc.add_table(rows=1, cols=max(1, len(header)))
    table.style = "Light Grid Accent 1"
    for i, cell_text in enumerate(header):
        cell = table.rows[0].cells[i]
        cell.text = ""
        run = cell.paragraphs[0].add_run(re.sub(r"\*\*?", "", cell_text))
        run.bold = True
    for row in rows:
        cells = table.add_row().cells
        for i in range(len(header)):
            value = row[i] if i < len(row) else ""
            cells[i].text = ""
            _add_runs(cells[i].paragraphs[0], value)


def markdown_to_docx(md_text, doc=None):
    """Convert proposal-grade markdown (headings, lists, tables, bold/italic)
    into a python-docx Document. Returns the Document."""
    doc = doc or Document()
    lines = (md_text or "").splitlines()
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        if not line.strip():
            i += 1
            continue
        # table block
        if line.lstrip().startswith("|") and i + 1 < len(lines) and _is_table_sep(lines[i + 1]):
            header = _split_row(line)
            i += 2
            rows = []
            while i < len(lines) and lines[i].lstrip().startswith("|"):
                rows.append(_split_row(lines[i]))
                i += 1
            _add_md_table(doc, header, rows)
            continue
        m = re.match(r"^(#{1,4})\s+(.*)$", line)
        if m:
            doc.add_heading(re.sub(r"\*\*?", "", m.group(2)).strip(), level=len(m.group(1)))
            i += 1
            continue
        m = re.match(r"^\s*[-*]\s+(.*)$", line)
        if m:
            p = doc.add_paragraph(style="List Bullet")
            _add_runs(p, m.group(1))
            i += 1
            continue
        m = re.match(r"^\s*\d+[.)]\s+(.*)$", line)
        if m:
            p = doc.add_paragraph(style="List Number")
            _add_runs(p, m.group(1))
            i += 1
            continue
        p = doc.add_paragraph()
        _add_runs(p, line.strip())
        i += 1
    return doc


def docx_bytes(doc):
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def narrative_docx(title, md_text, opp=None):
    """Standard proposal volume: header block + markdown body."""
    doc = Document()
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.LEFT
    if opp:
        meta = doc.add_paragraph()
        run = meta.add_run(
            f"{opp.get('title', '')} — {opp.get('solNumber') or 'TBD'} — "
            f"{opp.get('agency', '')}  |  Generated "
            f"{datetime.now(timezone.utc).strftime('%Y-%m-%d')}")
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x60, 0x60, 0x60)
    markdown_to_docx(md_text, doc)
    return docx_bytes(doc)


# ─────────────────────────── capability one-pager ────────────────────────────

def capability_docx(content, opp, org_name, rendering_png=None):
    doc = Document()
    doc.add_heading(content.get("title") or "Proposed Capability", level=0)
    meta = doc.add_paragraph()
    run = meta.add_run(
        f"{org_name} — {opp.get('solNumber') or 'TBD'} — {opp.get('agency', '')}  |  "
        f"{datetime.now(timezone.utc).strftime('%Y-%m-%d')}")
    run.font.size = Pt(9)

    doc.add_heading("Abstract", level=1)
    doc.add_paragraph(content.get("abstract") or "")

    doc.add_heading("Executive Summary", level=1)
    markdown_to_docx(content.get("executiveSummary") or "", doc)

    if content.get("keywords"):
        doc.add_heading("Keywords", level=1)
        doc.add_paragraph(", ".join(content["keywords"]))

    if rendering_png:
        doc.add_heading("Concept Rendering", level=1)
        try:
            doc.add_picture(io.BytesIO(rendering_png), width=Inches(6.2))
        except Exception:
            doc.add_paragraph("[Rendering could not be embedded]")

    sow = content.get("sow") or {}
    doc.add_heading("Statement of Work", level=1)
    if sow.get("scope"):
        doc.add_paragraph(sow["scope"])
    for t in sow.get("tasks") or []:
        doc.add_heading(f"Task {t.get('number', '')} — {t.get('title', '')}", level=2)
        doc.add_paragraph(t.get("description") or "")
        for d in t.get("deliverables") or []:
            doc.add_paragraph(d, style="List Bullet")

    wbs = content.get("wbs") or []
    if wbs:
        doc.add_heading("Work Breakdown Structure & Schedule", level=1)
        months = int(content.get("scheduleMonths") or 12)
        doc.add_paragraph(f"Period of performance: {months} months")
        _add_md_table(
            doc, ["WBS", "Task", "Owner", "Start (mo)", "End (mo)"],
            [[str(w.get("code", "")), str(w.get("task", "")), str(w.get("owner", "")),
              str(w.get("startMonth", "")), str(w.get("endMonth", ""))] for w in wbs])

    budget = content.get("budget") or {}
    items = budget.get("items") or []
    if items:
        doc.add_heading("Budget", level=1)
        total = sum(float(i.get("cost") or 0) for i in items)
        _add_md_table(
            doc, ["Category", "Description", "Cost (USD)"],
            [[str(i.get("category", "")), str(i.get("description", "")),
              f"{float(i.get('cost') or 0):,.0f}"] for i in items]
            + [["**Total**", "", f"**{total:,.0f}**"]])
        if budget.get("narrative"):
            doc.add_paragraph(budget["narrative"])

    for tbl in content.get("tables") or []:
        doc.add_heading(tbl.get("title") or "Table", level=1)
        _add_md_table(doc, [str(h) for h in (tbl.get("headers") or [])],
                      [[str(c) for c in row] for row in (tbl.get("rows") or [])])

    for chart in content.get("charts") or []:
        doc.add_heading(chart.get("title") or "Chart", level=1)
        _add_md_table(doc, ["Name", "Value"],
                      [[str(d.get("name", "")), str(d.get("value", ""))]
                       for d in (chart.get("data") or [])])

    return docx_bytes(doc)


# ─────────────────────────────── cost xlsx ───────────────────────────────────

HEADER_FILL = PatternFill("solid", fgColor="1F3864")
HEADER_FONT = Font(bold=True, color="FFFFFF")


def cost_volume_xlsx(content_json, opp, org_name):
    wb = Workbook()
    ws = wb.active
    ws.title = "Cost Volume"
    ws["A1"] = f"Cost / Budget Volume — {opp.get('title', '')}"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A2"] = (f"{org_name}  |  Sol # {opp.get('solNumber') or 'TBD'}  |  "
                f"{opp.get('agency', '')}  |  "
                f"{datetime.now(timezone.utc).strftime('%Y-%m-%d')}")
    ws["A2"].font = Font(size=9, italic=True)

    headers = ["Category", "Item", "Basis of Estimate", "Cost (USD)"]
    row0 = 4
    for col, h in enumerate(headers, start=1):
        c = ws.cell(row=row0, column=col, value=h)
        c.fill = HEADER_FILL
        c.font = HEADER_FONT
        c.alignment = Alignment(vertical="center")
    rows = content_json.get("rows") or []
    r = row0 + 1
    for item in rows:
        ws.cell(row=r, column=1, value=item.get("category", ""))
        ws.cell(row=r, column=2, value=item.get("item", ""))
        ws.cell(row=r, column=3, value=item.get("basis", ""))
        cost = ws.cell(row=r, column=4, value=float(item.get("cost") or 0))
        cost.number_format = '#,##0'
        r += 1
    total_label = ws.cell(row=r, column=3, value="TOTAL")
    total_label.font = Font(bold=True)
    total = ws.cell(row=r, column=4, value=f"=SUM(D{row0 + 1}:D{max(row0 + 1, r - 1)})")
    total.font = Font(bold=True)
    total.number_format = '#,##0'

    ws.column_dimensions["A"].width = 22
    ws.column_dimensions["B"].width = 38
    ws.column_dimensions["C"].width = 60
    ws.column_dimensions["D"].width = 16

    ws2 = wb.create_sheet("Narrative & Assumptions")
    ws2["A1"] = "Cost Narrative"
    ws2["A1"].font = Font(bold=True, size=12)
    ws2["A2"] = content_json.get("narrative", "")
    ws2["A2"].alignment = Alignment(wrap_text=True, vertical="top")
    ws2.row_dimensions[2].height = 120
    ws2["A4"] = "Assumptions"
    ws2["A4"].font = Font(bold=True, size=12)
    for i, a in enumerate(content_json.get("assumptions") or [], start=5):
        ws2.cell(row=i, column=1, value=f"• {a}")
    ws2.column_dimensions["A"].width = 110

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ─────────────────────────────── deck pptx ───────────────────────────────────

def briefing_pptx(content_json, opp, org_name):
    prs = Presentation()
    slides = content_json.get("slides") or []
    for idx, slide_def in enumerate(slides):
        layout = prs.slide_layouts[0 if idx == 0 else 1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_def.get("title", "")
        bullets = slide_def.get("bullets") or []
        if idx == 0:
            if len(slide.placeholders) > 1:
                sub = slide.placeholders[1]
                sub.text = (bullets[0] if bullets else
                            f"{org_name} — {opp.get('solNumber') or 'TBD'}")
        else:
            if len(slide.placeholders) > 1:
                body = slide.placeholders[1].text_frame
                body.clear()
                for i, b in enumerate(bullets):
                    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    p.text = str(b)
                    p.font.size = PptPt(18)
        notes = slide_def.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────── package zip ─────────────────────────────────

def package_zip(files):
    """files: list of (filename, bytes). Returns zip bytes."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data in files:
            zf.writestr(name, data)
    return buf.getvalue()
